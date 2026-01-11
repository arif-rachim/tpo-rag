#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Multi‑format ingest → chunk → embed → store (ChromaDB + BM25)

Supported extensions:
    *.pdf   – PyMuPDF
    *.docx  – python‑docx
    *.pptx  – python‑pptx
    *.xlsx / *.xls / *.xlsm – Excel (via the provided excel_converter)
"""

# --------------------------------------------------------------------------- #
#  Standard‑library imports
# --------------------------------------------------------------------------- #
import os
import pickle
import re
import warnings
from datetime import datetime, timezone
from functools import partial
from multiprocessing import Pool, cpu_count
from pathlib import Path
from typing import Any, Callable, Dict, List, Sequence, Tuple

# Import logging configuration
import sys
sys.path.insert(0, str(Path(__file__).parent.parent))
from logging_config import (
    get_ingestion_logger,
    log_exception,
    PerformanceLogger,
    log_system_info,
)

# --------------------------------------------------------------------------- #
#  Third‑party imports
# --------------------------------------------------------------------------- #
import chromadb
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from rank_bm25 import BM25Okapi
from sentence_transformers import SentenceTransformer
from tqdm import tqdm
from transformers import pipeline, AutoModelForTokenClassification, AutoTokenizer

# BiDi support for Arabic/Hebrew text
try:
    from bidi.algorithm import get_display
    BIDI_AVAILABLE = True
except ImportError:
    BIDI_AVAILABLE = False
    warnings.warn("python-bidi not installed. Arabic/Hebrew text may not display correctly.")

# --------------------------------------------------------------------------- #
#  Project imports
# --------------------------------------------------------------------------- #
import config
from .excel_converter import excel_to_json

# --------------------------------------------------------------------------- #
#  Global configuration & constants
# --------------------------------------------------------------------------- #
warnings.filterwarnings("ignore")

# Initialize logger
logger = get_ingestion_logger("ingestor.ingest")

INPUT_FOLDER = config.PATH_DOCUMENTS
VECTOR_DB_PATH = config.PATH_VECTOR_DB_STORAGE
BM25_INDEX_PATH = config.PATH_BM25_INDEX_FILE

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100
SEMANTIC_WEIGHT = 0.7

EMBEDDER_MODEL_PATH = config.PATH_MODEL_EMBEDDER
NER_MODEL_PATH = config.PATH_MODEL_NER
COLLECTION_NAME = config.VECTOR_DB_COLLECTION_NAME

NUM_WORKERS = min(cpu_count() - 1, 8) or 1

# Force HuggingFace to stay offline (keeps the original behaviour)
config.configure_offline_mode()

# Log configuration on module load
logger.info("Ingestion module loaded with configuration:")
logger.info(f"  Input folder: {INPUT_FOLDER}")
logger.info(f"  Vector DB path: {VECTOR_DB_PATH}")
logger.info(f"  BM25 index path: {BM25_INDEX_PATH}")
logger.info(f"  Chunk size: {CHUNK_SIZE}, overlap: {CHUNK_OVERLAP}")
logger.info(f"  Number of workers: {NUM_WORKERS}")
logger.info(f"  Embedder model: {EMBEDDER_MODEL_PATH}")
logger.info(f"  NER model: {NER_MODEL_PATH}")

# --------------------------------------------------------------------------- #
#  Regex patterns used for bilingual metadata extraction
# --------------------------------------------------------------------------- #
PATTERNS = {
    "jac_reg": r"JAC\s+REG\s+\d+-\d+",
    "jac_sgl": r"JAC\s+SGL\s+\d+-\d+\.\d+",
    "sop": r"\bSOP\b",
    "procedure": r"\b(?:Procedure|إجراء)\b",
}


# --------------------------------------------------------------------------- #
#  Helper utilities
# --------------------------------------------------------------------------- #
def detect_language(text: str) -> str:
    """Very cheap heuristic – Arabic if >50 Arabic chars in the first 500 chars."""
    arabic = len(re.findall(r"[\u0600-\u06FF]", text[:500]))
    return "ar" if arabic > 50 else "en"


def truncate_to_word_boundary(text: str, max_len: int, from_end: bool = False) -> str:
    """Trim to a word boundary without cutting a word."""
    if len(text) <= max_len:
        return text
    if from_end:
        tail = text[-max_len:]
        idx = tail.find(" ")
        return tail[idx + 1:] if idx > 0 else tail
    else:
        head = text[:max_len]
        idx = head.rfind(" ")
        return head[:idx] if idx > 0 else head


def chunk_text(text: str) -> List[str]:
    """Smart chunker that works for English & Arabic."""
    if not text or len(text) < 50:
        return []
    chunks, cur = [], ""
    paragraphs = re.split(r"\n\n+|(?<=[.!?؟।])\s+", text)

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(cur) + len(para) < CHUNK_SIZE:
            cur += para + "\n\n"
        else:
            if cur.strip():
                chunks.append(cur.strip())
            overlap = truncate_to_word_boundary(cur, CHUNK_OVERLAP, from_end=True)
            cur = overlap + para + "\n\n"

    if cur.strip():
        chunks.append(cur.strip())
    return chunks if chunks else [text[:CHUNK_SIZE]]


def _add_text_page(out: List[Dict], page_num: int, text: str) -> None:
    """Add a text page to output list."""
    # Fix BiDi text before adding to output
    fixed_text = fix_bidi_text(text.strip())
    out.append({
        "page": page_num,
        "text": fixed_text,
        "lang": detect_language(fixed_text),
    })


def extract_patterns(text: str) -> Dict[str, str]:
    """Return the first few matches for each regex pattern."""
    flags = re.IGNORECASE | re.UNICODE
    return {
        k: ", ".join(list(set(matches))[:5])
        for k, p in PATTERNS.items()
        if (matches := re.findall(p, text, flags))
    }


def fs_timestamp_to_iso(ts: float | None) -> str | None:
    """POSIX → ISO‑8601 (UTC)."""
    if ts is None:
        return None
    return datetime.fromtimestamp(ts, tz=timezone.utc).isoformat(timespec="seconds")


def fix_bidi_text(text: str) -> str:
    """
    Fix bidirectional text (Arabic, Hebrew, etc.) by properly reordering characters.

    This function detects if the text contains RTL (right-to-left) characters
    and applies the Unicode BiDi algorithm to fix character ordering issues
    that occur when extracting text from PDFs.

    Args:
        text: Raw text that may contain RTL characters in wrong order

    Returns:
        Text with properly ordered characters for display
    """
    if not BIDI_AVAILABLE or not text:
        return text

    # Check if text contains RTL characters (Arabic: 0x0600-0x06FF, Hebrew: 0x0590-0x05FF)
    has_rtl = any(
        '\u0590' <= char <= '\u05FF' or  # Hebrew
        '\u0600' <= char <= '\u06FF' or  # Arabic
        '\u0750' <= char <= '\u077F' or  # Arabic Supplement
        '\u08A0' <= char <= '\u08FF'     # Arabic Extended-A
        for char in text
    )

    if not has_rtl:
        return text

    try:
        # Apply BiDi algorithm line by line to preserve structure
        lines = text.split('\n')
        fixed_lines = []
        for line in lines:
            if line.strip():
                # get_display reshapes and reorders the text for proper display
                fixed_line = get_display(line)
                fixed_lines.append(fixed_line)
            else:
                fixed_lines.append(line)
        return '\n'.join(fixed_lines)
    except Exception as e:
        logger.warning(f"BiDi processing failed: {e}, returning original text")
        return text


# --------------------------------------------------------------------------- #
#  Extraction functions – one per file type (all return the same schema)
# --------------------------------------------------------------------------- #
def extract_pdf(pdf_path: Path) -> List[Dict[str, Any]]:
    """PDF → list[{'page','text','lang'}] (tables are marked as lang='table')."""
    try:
        logger.debug(f"Extracting PDF: {pdf_path.name}")
        doc = fitz.open(pdf_path)
        out: List[Dict[str, Any]] = []
        for page_num, page in enumerate(doc, 1):
            txt = page.get_text()
            if txt.strip():
                # Fix BiDi text (Arabic, Hebrew, etc.)
                txt = fix_bidi_text(txt)
                out.append(
                    {"page": page_num, "text": txt, "lang": detect_language(txt)}
                )
            # Tables (PyMuPDF ≥ 1.23.0)
            for table in page.find_tables():
                extracted = table.extract()
                if extracted:
                    table_txt = "\n".join(
                        [
                            " | ".join(str(cell) for cell in row if cell)
                            for row in extracted
                            if any(row)
                        ]
                    )
                    if table_txt.strip():
                        # Fix BiDi text in tables too
                        table_txt = fix_bidi_text(table_txt)
                        out.append(
                            {
                                "page": page_num,
                                "text": f"[TABLE]\n{table_txt}",
                                "lang": "table",
                            }
                        )
        doc.close()
        logger.debug(f"Extracted {len(out)} pages from PDF: {pdf_path.name}")
        return out
    except Exception as exc:
        log_exception(logger, exc, f"extract_pdf({pdf_path.name})")
        raise


def extract_docx(file_path: Path) -> List[Dict[str, Any]]:
    """DOCX → list[{'page','text','lang'}]; tables are separate pages."""
    doc = Document(file_path)
    out: List[Dict[str, Any]] = []
    page_counter = 1
    buffer = ""

    # Create a mapping of table elements to table objects for lookup
    table_map = {}
    for table in doc.tables:
        table_map[id(table._element)] = table

    # Process elements in order to preserve document structure
    for element in doc.element.body.iterchildren():
        # Paragraphs
        if element.tag.endswith("}p"):
            txt = element.text or ""
            buffer += txt + "\n"
        # Tables
        elif element.tag.endswith("}tbl"):
            # Find the corresponding table object
            table_obj = table_map.get(id(element))
            if table_obj:
                rows = []
                for row in table_obj.rows:
                    cells = [cell.text or "" for cell in row.cells]
                    rows.append(" | ".join(cells))
                table_txt = "\n".join(rows)

                if buffer:
                    _add_text_page(out, page_counter, buffer)
                    buffer = ""
                    page_counter += 1
                # Fix BiDi text in tables
                table_txt = fix_bidi_text(table_txt)
                out.append(
                    {"page": page_counter, "text": f"[TABLE]\n{table_txt}", "lang": "table"}
                )
                page_counter += 1

    if buffer:
        _add_text_page(out, page_counter, buffer)
    return out


def extract_pptx(file_path: Path) -> List[Dict[str, Any]]:
    """PPTX → list[{'pagetext', 'lang'}]; tables are inline."""
    prs = Presentation(file_path)
    out: List[Dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_txt = []
        # Text boxes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_txt.append(shape.text)
        # Tables
        for shape in slide.shapes:
            if shape.has_table:
                rows = []
                for r in shape.table.rows:
                    cells = [c.text for c in r.cells]
                    rows.append(" | ".join(cells))
                table_txt = "\n".join(rows)
                slide_txt.append(f"[TABLE]\n{table_txt}")

        combined = "\n".join(slide_txt).strip()
        if combined:
            # Fix BiDi text in slides
            combined = fix_bidi_text(combined)
            out.append(
                {"page": idx, "text": combined, "lang": detect_language(combined)}
            )
    return out


def extract_excel(file_path: Path) -> List[Dict[str, Any]]:
    """
    Excel → list[{'page','text','lang','sheet_title','total_cells'}].

    * ``page`` is the **sheet title** (string) – this makes the
      identifier human readable.
    * ``sheet_title`` is kept for backward‑compatibility with any
      downstream code that expects the original camel‑case name.
    """
    workbook_dict = excel_to_json(str(file_path), sample_size=None)
    if not workbook_dict:
        return []  # conversion failed – caller will treat as empty

    pages: List[Dict[str, Any]] = []
    for sheet_index, sheet in enumerate(workbook_dict.get("sheets", []), start=1):
        # Prefer the explicit title; fall back to a generic one.
        sheet_title = sheet.get("sheetTitle", f"Sheet{sheet_index}")
        cells = sheet.get("cells", {})

        # Build a readable text block – one line per cell.
        lines: List[str] = []
        for coord, cell in sorted(cells.items()):
            # Prefer the calculated value if it exists, otherwise raw value.
            value = cell.get("calculated_value") or cell.get("value")
            if value is None:
                continue
            formula = f"  (formula: {cell['formula']})" if "formula" in cell else ""
            lines.append(f"{coord}: {value}{formula}")

        sheet_text = "\n".join(lines)
        if not sheet_text:
            # Skip empty sheets – they add noise.
            continue

        pages.append(
            {
                # NOTE: ``page`` is now the *sheet title* (string) – not an int.
                "page": sheet_title,
                "text": sheet_text,
                "lang": detect_language(sheet_text),
                # New snake_case name for the rest of the pipeline.
                "sheet_title": sheet_title,
                "total_cells": len(cells),
            }
        )
    return pages


# --------------------------------------------------------------------------- #
#  Dispatcher – maps suffix → extractor
# --------------------------------------------------------------------------- #
EXTRACTOR_MAP: Dict[str, Callable[[Path], List[Dict[str, Any]]]] = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_excel,
    ".xls": extract_excel,
    ".xlsm": extract_excel,
    # Extend here if you later want .txt, .doc, .ppt, etc.
}


def extract_file(file_path: Path) -> List[Dict[str, Any]]:
    """Pick the right extractor based on the file suffix."""
    suffix = file_path.suffix.lower()
    if suffix not in EXTRACTOR_MAP:
        raise ValueError(f"Unsupported file type: {suffix}")
    return EXTRACTOR_MAP[suffix](file_path)


# --------------------------------------------------------------------------- #
#  Embedding / storage helpers (unchanged)
# --------------------------------------------------------------------------- #
def create_embeddings(texts: Sequence[str], model: SentenceTransformer):
    prefixed = [f"passage: {t[:512]}" for t in texts]
    return model.encode(
        prefixed,
        show_progress_bar=False,
        batch_size=64,
        normalize_embeddings=True,
    )


def store_chromadb(
        chunks: Sequence[str],
        embeddings,
        metadata_list: Sequence[Dict[str, Any]],
        client: chromadb.PersistentClient,
):
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        metadata={"description": "Arabic‑English bilingual technical documents"},
    )
    ids = [
        f"{meta.get('filename', 'doc')}_p{meta.get('page', 0)}_c{i}"
        for i, meta in enumerate(metadata_list)
    ]

    batch = 100
    for i in range(0, len(chunks), batch):
        end = min(i + batch, len(chunks))
        collection.add(
            documents=chunks[i:end],
            embeddings=embeddings[i:end].tolist()
            if hasattr(embeddings, "tolist")
            else embeddings[i:end],
            metadatas=metadata_list[i:end],
            ids=ids[i:end],
        )
    return collection


def build_bm25(chunks: Sequence[str], metadatas: Sequence[Dict[str, Any]]):
    tokenized = [re.findall(r"\w+", c.lower()) for c in chunks]
    bm25 = BM25Okapi(tokenized)
    with open(BM25_INDEX_PATH, "wb") as f:
        pickle.dump({"bm25": bm25, "chunks": chunks, "metadatas": metadatas}, f)
    return bm25


# --------------------------------------------------------------------------- #
#  Worker that runs in a separate process (extract → chunk → meta)
# --------------------------------------------------------------------------- #
def extract_and_chunk_file(
        file_path: Path,
        ner: Any = None,
        input_folder: Path = None,
) -> Tuple[str, List[str], List[Dict[str, Any]]] | None:
    """
    Returns (filename, list_of_chunks, list_of_metadatas) or None on failure.
    """
    try:
        # Compute relative path from input folder (or use basename if not provided)
        if input_folder:
            try:
                relative_path = file_path.relative_to(input_folder)
                filename = str(relative_path).replace('\\', '/')  # Normalize path separators
            except ValueError:
                # If file_path is not relative to input_folder, use basename
                filename = file_path.name
        else:
            filename = file_path.name

        created_iso = fs_timestamp_to_iso(os.path.getctime(file_path))
        modified_iso = fs_timestamp_to_iso(os.path.getmtime(file_path))
        file_size = file_path.stat().st_size  # Get file size in bytes

        # --------------------------------------------------------------- #
        #  1️⃣  Extract raw pages (text + optional tables/formulas)
        # --------------------------------------------------------------- #
        pages = extract_file(file_path)
        if not pages:
            return None

        # ``total_pages`` is now simply the number of extracted pages/sheets.
        total_pages = len(pages)

        # --------------------------------------------------------------- #
        #  2️⃣  Chunk & enrich metadata
        # --------------------------------------------------------------- #
        all_chunks, all_meta = [], []
        for page in pages:
            chunks = chunk_text(page["text"])
            for idx, chunk in enumerate(chunks):
                meta = {
                    **extract_patterns(chunk),
                    "filename": filename,
                    "page": page["page"],
                    "total_pages": total_pages,
                    "chunk": idx,
                    "lang": page.get("lang", "unknown"),
                    "created_at": created_iso,
                    "modified_at": modified_iso,
                    "file_size": file_size,
                    "file_type": file_path.suffix.lower().lstrip("."),
                    **{k: page[k] for k in ("sheet_title", "total_cells") if k in page}
                }

                # ---- optional NER (first chunk of each page) ---- #
                if ner and idx == 0 and len(chunk) > 100:
                    try:
                        ents = ner(chunk[:500], aggregation_strategy="simple")
                        for entity_type, key in [("PER", "persons"), ("ORG", "orgs")]:
                            entities = [e["word"] for e in ents if entity_type in e["entity_group"]]
                            if entities:
                                meta[key] = ", ".join(set(entities)[:3])
                    except Exception:
                        pass  # NER errors should never abort the pipeline

                all_chunks.append(chunk)
                all_meta.append(meta)

        return filename, all_chunks, all_meta

    except Exception as exc:
        # Log the full error with traceback
        log_exception(logger, exc, f"extract_and_chunk_file({file_path.name})")
        logger.error(f"Failed processing file: {file_path.name}")
        return None


# --------------------------------------------------------------------------- #
#  Main orchestration
# --------------------------------------------------------------------------- #
def main() -> None:
    logger.info("=" * 80)
    logger.info("Document Ingestion Process Starting")
    logger.info("=" * 80)
    log_system_info(logger)

    with PerformanceLogger(logger, "Complete ingestion process"):
        # --------------------------------------------------------------- #
        #  1️⃣  Find *all* supported documents
        # --------------------------------------------------------------- #
        logger.info(f"Scanning for documents in: {INPUT_FOLDER}")
        all_files: List[Path] = []
        for suffix in EXTRACTOR_MAP.keys():
            files = list(Path(INPUT_FOLDER).rglob(f"*{suffix}"))
            if files:
                logger.info(f"Found {len(files)} {suffix} file(s)")
            all_files.extend(files)

        if not all_files:
            logger.warning("No supported documents found - exiting")
            logger.info("Supported file types: " + ", ".join(EXTRACTOR_MAP.keys()))
            return

        logger.info(f"Total files to process: {len(all_files)}")

        # --------------------------------------------------------------- #
        #  2️⃣  Load heavy models once (in the main process)
        # --------------------------------------------------------------- #
        logger.info("Loading embedding model...")
        with PerformanceLogger(logger, "Loading embedder model"):
            embedder = SentenceTransformer(
                str(EMBEDDER_MODEL_PATH),
                device="cpu",
                local_files_only=True,
                tokenizer_kwargs={"clean_up_tokenization_spaces": True, "fix_mistral_regex": True},
            )
        logger.info("Embedding model loaded successfully")

        logger.info("Initializing ChromaDB client...")
        with PerformanceLogger(logger, "Initializing ChromaDB"):
            client = chromadb.PersistentClient(path=VECTOR_DB_PATH)
        logger.info("ChromaDB client initialized")

        # Clear existing collection and BM25 index to sync with filesystem
        # This ensures deleted files are removed from the database
        logger.info("Clearing existing ChromaDB collection to sync with filesystem...")
        try:
            client.delete_collection(name=COLLECTION_NAME)
            logger.info(f"Successfully deleted old collection: {COLLECTION_NAME}")
        except Exception as e:
            logger.info(f"No existing collection to delete (this is normal for first run): {e}")

        # Clear BM25 index file
        logger.info("Clearing existing BM25 index...")
        try:
            if BM25_INDEX_PATH.exists():
                BM25_INDEX_PATH.unlink()
                logger.info(f"Successfully deleted old BM25 index: {BM25_INDEX_PATH}")
            else:
                logger.info("No existing BM25 index to delete (this is normal for first run)")
        except Exception as e:
            logger.warning(f"Could not delete BM25 index: {e}")

        logger.info("ChromaDB and BM25 index are now synced with filesystem - ready for fresh ingestion")

        logger.info("Loading NER model...")
        with PerformanceLogger(logger, "Loading NER model"):
            ner_pipe = pipeline(
                "ner",
                model=AutoModelForTokenClassification.from_pretrained(
                    str(NER_MODEL_PATH), local_files_only=True
                ),
                tokenizer=AutoTokenizer.from_pretrained(
                    str(NER_MODEL_PATH), local_files_only=True, clean_up_tokenization_spaces=True
                ),
                device=-1,  # CPU
                aggregation_strategy="simple",
            )
        logger.info("NER model loaded successfully")

        worker = partial(extract_and_chunk_file, ner=ner_pipe, input_folder=Path(INPUT_FOLDER))

        # --------------------------------------------------------------- #
        #  3️⃣  Parallel extraction / chunking (CPU bound)
        # --------------------------------------------------------------- #
        logger.info(f"Starting parallel processing with {NUM_WORKERS} workers")
        total_chunks = processed = failed = 0

        with PerformanceLogger(logger, "Parallel file processing"):
            with Pool(processes=NUM_WORKERS) as pool, tqdm(
                    total=len(all_files), desc="Processing files", unit="file"
            ) as pbar:
                for result in pool.imap_unordered(worker, all_files):
                    if result is None:
                        failed += 1
                        pbar.update(1)
                        continue

                    filename, chunks, metadata = result
                    processed += 1
                    logger.debug(f"Successfully processed: {filename} ({len(chunks)} chunks)")

                    if chunks:
                        try:
                            with PerformanceLogger(logger, f"Embedding and storing {filename}"):
                                embeddings = create_embeddings(chunks, embedder)
                                store_chromadb(chunks, embeddings, metadata, client)
                                total_chunks += len(chunks)
                            logger.debug(f"Stored {len(chunks)} chunks for {filename}")
                        except Exception as exc:
                            tqdm.write(f"❌  Storing error for {filename}: {exc}")
                            log_exception(logger, exc, f"storing {filename}")
                            failed += 1

                    pbar.update(1)
                    pbar.set_postfix(
                        {"processed": processed, "failed": failed, "chunks": total_chunks}
                    )

        logger.info(f"File processing complete: {processed} processed, {failed} failed")

        # --------------------------------------------------------------- #
        #  4️⃣  Build BM25 index from everything that landed in Chroma
        # --------------------------------------------------------------- #
        if total_chunks:
            logger.info("Building BM25 index...")
            with PerformanceLogger(logger, "Building BM25 index"):
                collection = client.get_collection(COLLECTION_NAME)
                docs = collection.get()
                if docs and docs["documents"]:
                    build_bm25(docs["documents"], docs["metadatas"])
                    logger.info(f"BM25 index built with {len(docs['documents'])} documents")
                else:
                    logger.warning("No documents found in ChromaDB collection")
        else:
            logger.warning("No chunks to index - skipping BM25 build")

    logger.info("=" * 80)
    logger.info("Document Ingestion Complete")
    logger.info("=" * 80)
    logger.info(f"Summary:")
    logger.info(f"  Total files found: {len(all_files)}")
    logger.info(f"  Successfully processed: {processed}")
    logger.info(f"  Failed: {failed}")
    logger.info(f"  Total chunks indexed: {total_chunks}")
    logger.info(f"  Success rate: {(processed / len(all_files) * 100):.1f}%")
    logger.info("=" * 80)


if __name__ == "__main__":
    main()
